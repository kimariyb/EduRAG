from __future__ import annotations

from collections.abc import Iterator
from io import BytesIO
from pathlib import Path
from typing import Any

import fitz
import numpy as np
from docx import Document as open_docx
from docx.document import Document as DocxDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.parts.image import ImagePart
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tqdm import tqdm

from base.logger import logger


log = logger.bind(module=__name__)


def get_ocr(use_cuda: bool = True) -> Any:
    try:
        from rapidocr_paddle import RapidOCR

        return RapidOCR(
            det_use_cuda=use_cuda,
            cls_use_cuda=use_cuda,
            rec_use_cuda=use_cuda,
        )
    except ImportError:
        # Fall back to ONNX Runtime when Paddle is unavailable on CPU systems.
        from rapidocr_onnxruntime import RapidOCR

        log.info(
            "rapidocr_paddle is unavailable; falling back to "
            "rapidocr_onnxruntime"
        )
        return RapidOCR()


def _ocr_lines(ocr: Any, image: Any) -> list[str]:
    result, _ = ocr(image)
    if not result:
        return []
    return [str(line[1]) for line in result]


def _image_bytes_to_array(blob: bytes) -> np.ndarray:
    with Image.open(BytesIO(blob)) as image:
        return np.asarray(image.convert("RGB"))


class PlainTextLoader(BaseLoader):
    def __init__(self, file_path: str, encoding: str = "utf-8") -> None:
        self.file_path = file_path
        self.encoding = encoding

    def lazy_load(self) -> Iterator[Document]:
        text = Path(self.file_path).read_text(encoding=self.encoding)
        log.info(
            "Loaded plain text document: path={}, characters={}",
            self.file_path,
            len(text),
        )
        yield Document(
            page_content=text,
            metadata={"source": self.file_path},
        )


class PDFLoader(BaseLoader):
    def __init__(self, file_path: str) -> None:
        self.file_path = file_path

    def lazy_load(self) -> Iterator[Document]:
        text = self.pdf2text()
        yield Document(page_content=text, metadata={"source": self.file_path})

    def pdf2text(self) -> str:
        ocr = get_ocr()
        document = fitz.open(self.file_path)
        parts: list[str] = []
        log.info(
            "Loading PDF document: path={}, pages={}",
            self.file_path,
            document.page_count,
        )

        try:
            with tqdm(
                total=document.page_count,
                desc=f"OCR PDF page: 0/{document.page_count}",
            ) as progress:
                for index, page in enumerate(document):
                    progress.set_description(
                        f"OCR PDF page: {index + 1}/{document.page_count}"
                    )

                    text = page.get_text()
                    if text:
                        parts.append(text.rstrip("\n"))

                    for image_info in page.get_image_info(xrefs=True):
                        xref = image_info.get("xref")
                        if not xref:
                            continue

                        pixmap = fitz.Pixmap(document, xref)
                        # OCR expects three-channel RGB input.
                        if pixmap.n != 3:
                            pixmap = fitz.Pixmap(fitz.csRGB, pixmap)
                        image_array = np.frombuffer(
                            pixmap.samples,
                            dtype=np.uint8,
                        ).reshape(pixmap.height, pixmap.width, 3)

                        rotation = int(page.rotation) % 360
                        if rotation:
                            # PDF rotation is clockwise; np.rot90 is counterclockwise.
                            image_array = np.ascontiguousarray(
                                np.rot90(
                                    image_array,
                                    k=-(rotation // 90),
                                )
                            )
                        parts.extend(_ocr_lines(ocr, image_array))
                    progress.update(1)
        finally:
            document.close()

        text = "\n".join(parts)
        log.info(
            "Loaded PDF document: path={}, characters={}",
            self.file_path,
            len(text),
        )
        return text


class DOCXLoader(BaseLoader):
    def __init__(self, filepath: str) -> None:
        self.filepath = filepath

    def lazy_load(self) -> Iterator[Document]:
        text = self.doc2text(self.filepath)
        yield Document(page_content=text, metadata={"source": self.filepath})

    def doc2text(self, filepath: str) -> str:
        ocr = get_ocr()
        document = open_docx(filepath)
        parts: list[str] = []
        total_blocks = len(document.paragraphs) + len(document.tables)
        log.info(
            "Loading DOCX document: path={}, blocks={}",
            filepath,
            total_blocks,
        )

        with tqdm(
            total=total_blocks,
            desc=f"OCR DOCX block: 0/{total_blocks}",
        ) as progress:
            for index, block in enumerate(self._iter_block_items(document)):
                progress.set_description(
                    f"OCR DOCX block: {index + 1}/{total_blocks}"
                )

                if isinstance(block, Paragraph):
                    parts.append(block.text.strip())
                    for image in block._element.xpath(".//pic:pic"):
                        for image_id in image.xpath(".//a:blip/@r:embed"):
                            part = document.part.related_parts[image_id]
                            if isinstance(part, ImagePart):
                                image_array = _image_bytes_to_array(part.blob)
                                parts.extend(_ocr_lines(ocr, image_array))
                elif isinstance(block, Table):
                    for row in block.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                parts.append(paragraph.text.strip())
                progress.update(1)

        text = "\n".join(part for part in parts if part)
        log.info(
            "Loaded DOCX document: path={}, characters={}",
            filepath,
            len(text),
        )
        return text

    @staticmethod
    def _iter_block_items(
        parent: DocxDocument | _Cell,
    ) -> Iterator[Paragraph | Table]:
        # Iterate over XML children to preserve paragraph and table order.
        if isinstance(parent, DocxDocument):
            parent_element = parent.element.body
        elif isinstance(parent, _Cell):
            parent_element = parent._tc
        else:
            parent_type = type(parent).__name__
            raise ValueError(f"Unsupported DOCX block parent: {parent_type}")

        for child in parent_element.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)


class PPTXLoader(BaseLoader):
    def __init__(self, filepath: str) -> None:
        self.filepath = filepath

    def lazy_load(self) -> Iterator[Document]:
        text = self.ppt2text(self.filepath)
        yield Document(page_content=text, metadata={"source": self.filepath})

    def ppt2text(self, filepath: str) -> str:
        presentation = Presentation(filepath)
        ocr = get_ocr()
        parts: list[str] = []
        log.info(
            "Loading PPTX document: path={}, slides={}",
            filepath,
            len(presentation.slides),
        )

        def extract_text(shape: Any) -> None:
            if shape.has_text_frame:
                parts.append(shape.text.strip())

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            parts.append(paragraph.text.strip())

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_array = _image_bytes_to_array(shape.image.blob)
                parts.extend(_ocr_lines(ocr, image_array))
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child_shape in shape.shapes:
                    extract_text(child_shape)

        with tqdm(
            total=len(presentation.slides),
            desc=f"OCR PPTX slide: 0/{len(presentation.slides)}",
        ) as progress:
            for slide_number, slide in enumerate(presentation.slides, start=1):
                progress.set_description(
                    f"OCR PPTX slide: {slide_number}/{len(presentation.slides)}"
                )
                # Extract shapes in top-to-bottom, left-to-right reading order.
                sorted_shapes = sorted(
                    slide.shapes,
                    key=lambda shape: (shape.top, shape.left),
                )
                for shape in sorted_shapes:
                    extract_text(shape)
                progress.update(1)

        text = "\n".join(part for part in parts if part)
        log.info(
            "Loaded PPTX document: path={}, characters={}",
            filepath,
            len(text),
        )
        return text


class IMGLoader(BaseLoader):
    def __init__(self, img_path: str) -> None:
        self.img_path = img_path

    def lazy_load(self) -> Iterator[Document]:
        text = self.img2text()
        yield Document(page_content=text, metadata={"source": self.img_path})

    def img2text(self) -> str:
        ocr = get_ocr()
        text = "\n".join(_ocr_lines(ocr, self.img_path))
        log.info(
            "Loaded image document: path={}, characters={}",
            self.img_path,
            len(text),
        )
        return text
